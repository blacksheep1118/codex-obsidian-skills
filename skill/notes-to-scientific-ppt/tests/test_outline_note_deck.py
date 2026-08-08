from __future__ import annotations

import os
from pathlib import Path
import subprocess
import sys
from zipfile import ZipFile

import pytest
from pptx import Presentation

from scripts.build_scientific_deck import build_deck, load_or_create_brief
from scripts.outline_note_deck import build_vault_index, collect_linked_markdown_files, is_excluded


ROOT = Path(__file__).resolve().parents[1]


def run_script(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, "scripts/outline_note_deck.py", *args],
        cwd=ROOT,
        text=True,
        encoding="utf-8",
        errors="replace",
        capture_output=True,
        check=True,
    )


def run_script_unchecked(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, "scripts/outline_note_deck.py", *args],
        cwd=ROOT,
        text=True,
        encoding="utf-8",
        errors="replace",
        capture_output=True,
        check=False,
    )


def run_build_script(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, "scripts/build_scientific_deck.py", *args],
        cwd=ROOT,
        text=True,
        encoding="utf-8",
        errors="replace",
        capture_output=True,
        check=True,
    )


def run_build_script_unchecked(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, "scripts/build_scientific_deck.py", *args],
        cwd=ROOT,
        text=True,
        encoding="utf-8",
        errors="replace",
        capture_output=True,
        check=False,
    )


def slide_title(slide) -> str:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            return shape.text.strip().splitlines()[0]
    return ""


def test_outline_note_deck_creates_scientific_brief(tmp_path: Path):
    notes = tmp_path / "notes"
    notes.mkdir()
    note = notes / "01_method.md"
    note.write_text(
        "\n".join(
            [
                "# From Noise Modeling to Blind Denoising",
                "",
                "## 问题背景",
                "",
                "盲去噪需要在未知噪声下恢复图像。",
                "",
                "## 关键公式",
                "",
                "$$",
                "y = x + n",
                "$$",
                "",
                "## 实验",
                "",
                "| 数据集 | 指标 |",
                "| --- | --- |",
                "| BSD68 | PSNR |",
                "",
                "参考: [paper](https://example.com/paper.pdf)",
            ]
        ),
        encoding="utf-8",
    )
    out = tmp_path / "deck_brief.md"

    result = run_script(str(notes), "--out", str(out), "--title", "Blind Denoising", "--language", "en")

    assert f"wrote_deck_brief {out}" in result.stdout
    text = out.read_text(encoding="utf-8")
    assert "# Blind Denoising" in text
    assert "## Source Inventory" in text
    assert "## Extracted Note Structure" in text
    assert "## Evidence Ledger" in text
    assert "## Suggested Scientific Deck Spine" in text
    assert "Deck Mode: paper-reading" in text
    assert "## Draft Slide Backlog" in text
    assert "## Coverage Checklist" in text
    assert "科研严谨风" in text
    assert "问题背景" in text
    assert "关键公式" in text
    assert "equation-to-intuition bridge" in text
    assert "result/comparison table" in text
    assert "[formula/algorithm] Turn `关键公式`" in text
    assert "https://example.com/paper.pdf" in text


def test_outline_note_deck_creates_chinese_brief_for_chinese_notes(tmp_path: Path):
    note = tmp_path / "中文笔记.md"
    note.write_text(
        "# 盲图像去噪\n\n## 问题背景\n\n盲去噪需要处理未知噪声。\n\n## 局限\n\n真实噪声仍然复杂。\n",
        encoding="utf-8",
    )

    result = run_script(str(note), "--language", "zh")

    assert "## 来源盘点" in result.stdout
    assert "## 建议科学演示主线" in result.stdout
    assert "演示模式:" in result.stdout
    assert "## 覆盖检查清单" in result.stdout


def test_outline_note_deck_counts_wiki_embeds_and_cjk_chars(tmp_path: Path):
    note = tmp_path / "embed.md"
    note.write_text(
        "# Embed Note\n\n中文内容ABC words。\n\n![[figure 1.png]]\n\n![[dataset.csv]]\n",
        encoding="utf-8",
    )

    result = run_script(str(note), "--language", "en")

    assert "| `embed.md` | Embed Note |" in result.stdout
    assert " | 1 | 1 | 0 | 0 |" in result.stdout


def test_outline_note_deck_respects_max_slides(tmp_path: Path):
    note = tmp_path / "proposal.md"
    note.write_text("# Proposal\n\n## Method\n\nText.\n", encoding="utf-8")

    result = run_script(str(note), "--max-slides", "4", "--language", "en")
    spine = result.stdout.split("## Suggested Scientific Deck Spine", 1)[1].split("## Draft Slide Backlog", 1)[0]
    numbered = [line for line in spine.splitlines() if line.strip() and line.lstrip()[0].isdigit()]

    assert len(numbered) == 3
    assert "Maximum total slide count: 4 (including the title slide)" in result.stdout


@pytest.mark.parametrize("value", ["0", "-1"])
def test_outline_and_build_reject_nonpositive_max_slides(tmp_path: Path, value: str):
    note = tmp_path / "note.md"
    note.write_text("# Note\n\nEvidence.\n", encoding="utf-8")

    outline = run_script_unchecked(str(note), "--max-slides", value)
    build = run_build_script_unchecked(
        str(note),
        "--out",
        str(tmp_path / "deck.pptx"),
        "--max-slides",
        value,
    )

    assert outline.returncode != 0
    assert build.returncode != 0
    assert "at least 2" in outline.stderr
    assert "at least 2" in build.stderr


def test_programmatic_build_rejects_invalid_total_slide_budget(tmp_path: Path):
    note = tmp_path / "note.md"
    note.write_text("# Note\n\nEvidence.\n", encoding="utf-8")

    with pytest.raises(ValueError, match="at least 2"):
        build_deck(note, tmp_path / "deck.pptx", max_slides=0)


def test_build_from_brief_inherits_total_slide_budget_when_flag_is_omitted(tmp_path: Path):
    note = tmp_path / "note.md"
    brief = tmp_path / "brief.md"
    deck = tmp_path / "deck.pptx"
    note.write_text("# Note\n\n## Method\n\nGrounded evidence.\n", encoding="utf-8")
    run_script(str(note), "--out", str(brief), "--max-slides", "4", "--language", "en")

    result = run_build_script(str(brief), "--out", str(deck))

    assert "slides 4" in result.stdout
    assert len(Presentation(str(deck)).slides) == 4


def test_outline_note_deck_can_follow_local_wiki_links(tmp_path: Path):
    vault = tmp_path / "vault"
    main = vault / "main.md"
    linked = vault / "Linked.md"
    main.parent.mkdir()
    main.write_text("# Main\n\nSee [[Linked]].\n", encoding="utf-8")
    linked.write_text("# Linked\n\n## 实验\n\n| A | B |\n| --- | --- |\n| 1 | 2 |\n", encoding="utf-8")

    result = run_script(str(main), "--follow-links", "--vault-root", str(vault), "--max-depth", "1", "--language", "en")

    assert "`main.md`" in result.stdout
    assert "`Linked.md`" in result.stdout
    assert "result/comparison table" in result.stdout


def test_outline_note_deck_excludes_output_inside_scanned_notes_and_is_byte_stable(tmp_path: Path):
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "method.md").write_text("# Method\n\n## Evidence\n\nSource-grounded result.\n", encoding="utf-8")
    out = notes / "generated_deck_brief.md"

    run_script(str(notes), "--out", str(out), "--title", "Stable Brief", "--language", "en")
    first = out.read_bytes()
    run_script(str(notes), "--out", str(out), "--title", "Stable Brief", "--language", "en")

    assert out.read_bytes() == first
    assert "generated_deck_brief.md" not in out.read_text(encoding="utf-8")


def test_outline_note_deck_excludes_case_alias_for_existing_output_on_case_insensitive_fs(
    tmp_path: Path,
):
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "method.md").write_text("# Method\n\n## Evidence\n\nSource-grounded result.\n", encoding="utf-8")
    lower_out = notes / "generated_deck_brief.md"
    upper_out = notes / "GENERATED_DECK_BRIEF.md"

    run_script(str(notes), "--out", str(lower_out), "--title", "Stable Brief", "--language", "en")
    try:
        aliases_same_file = upper_out.exists() and upper_out.samefile(lower_out)
    except OSError:
        aliases_same_file = False
    if not aliases_same_file:
        pytest.skip("filesystem is case-sensitive")

    first = lower_out.read_bytes()
    run_script(str(notes), "--out", str(upper_out), "--title", "Stable Brief", "--language", "en")
    assert lower_out.read_bytes() == first
    run_script(str(notes), "--out", str(upper_out), "--title", "Stable Brief", "--language", "en")

    assert lower_out.read_bytes() == first


def test_output_exclusion_does_not_fold_distinct_case_paths_on_case_sensitive_fs(tmp_path: Path):
    lower = tmp_path / "generated_deck_brief.md"
    upper = tmp_path / "GENERATED_DECK_BRIEF.md"
    lower.write_text("lower\n", encoding="utf-8")
    upper.write_text("upper\n", encoding="utf-8")
    try:
        paths_are_distinct = not lower.samefile(upper)
    except OSError:
        paths_are_distinct = False
    if not paths_are_distinct:
        pytest.skip("filesystem is case-insensitive")

    assert not is_excluded(lower, {upper.resolve()})


def test_outline_note_deck_excludes_existing_hardlink_output_alias(tmp_path: Path):
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "method.md").write_text("# Method\n\n## Evidence\n\nSource-grounded result.\n", encoding="utf-8")
    original_out = notes / "generated_deck_brief.md"
    alias_out = notes / "brief_alias.md"

    run_script(str(notes), "--out", str(original_out), "--title", "Stable Brief", "--language", "en")
    try:
        os.link(original_out, alias_out)
    except OSError as exc:
        pytest.skip(f"hard links unavailable: {exc}")

    first = original_out.read_bytes()
    run_script(str(notes), "--out", str(alias_out), "--title", "Stable Brief", "--language", "en")
    assert original_out.read_bytes() == first
    run_script(str(notes), "--out", str(alias_out), "--title", "Stable Brief", "--language", "en")

    assert original_out.read_bytes() == first


def test_outline_note_deck_rejects_existing_symlink_output_alias(tmp_path: Path):
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "method.md").write_text("# Method\n\n## Evidence\n\nSource-grounded result.\n", encoding="utf-8")
    original_out = notes / "generated_deck_brief.md"
    alias_out = notes / "brief_alias.md"

    run_script(str(notes), "--out", str(original_out), "--title", "Stable Brief", "--language", "en")
    try:
        alias_out.symlink_to(original_out.name)
    except OSError as exc:
        pytest.skip(f"symbolic links unavailable: {exc}")

    first = original_out.read_bytes()
    result = run_script_unchecked(
        str(notes),
        "--out",
        str(alias_out),
        "--title",
        "Stable Brief",
        "--language",
        "en",
    )

    assert result.returncode == 1
    assert "symlink" in result.stderr.lower()
    assert original_out.read_bytes() == first


def test_outline_note_deck_rejects_dangling_symlink_output(tmp_path: Path):
    note = tmp_path / "note.md"
    note.write_text("# Note\n\nReal evidence.\n", encoding="utf-8")
    missing_target = tmp_path / "missing.md"
    out = tmp_path / "brief.md"
    out.symlink_to(missing_target)

    result = run_script_unchecked(str(note), "--out", str(out), "--language", "en")

    assert result.returncode == 1
    assert "symlink" in result.stderr.lower()
    assert out.is_symlink()
    assert not missing_target.exists()


def test_outline_note_deck_writes_regular_output_file(tmp_path: Path):
    note = tmp_path / "note.md"
    note.write_text("# Note\n\nReal evidence.\n", encoding="utf-8")
    out = tmp_path / "brief.md"

    result = run_script(str(note), "--out", str(out), "--language", "en")

    assert result.returncode == 0
    assert out.is_file()
    assert not out.is_symlink()
    assert "## Source Inventory" in out.read_text(encoding="utf-8")


def test_outline_note_deck_follow_links_excludes_linked_output_and_is_byte_stable(tmp_path: Path):
    vault = tmp_path / "vault"
    vault.mkdir()
    main = vault / "main.md"
    out = vault / "generated_deck_brief.md"
    main.write_text(
        "# Main\n\nSee [[Evidence]] and [[generated_deck_brief]].\n",
        encoding="utf-8",
    )
    (vault / "Evidence.md").write_text("# Evidence\n\n## Result\n\nMeasured evidence.\n", encoding="utf-8")
    args = (
        str(main),
        "--out",
        str(out),
        "--title",
        "Linked Stable Brief",
        "--language",
        "en",
        "--follow-links",
        "--vault-root",
        str(vault),
        "--max-depth",
        "1",
    )

    run_script(*args)
    first = out.read_bytes()
    run_script(*args)

    assert out.read_bytes() == first
    assert "generated_deck_brief.md" not in out.read_text(encoding="utf-8")


def test_outline_note_deck_does_not_follow_wiki_links_outside_vault(tmp_path: Path):
    vault = tmp_path / "vault"
    vault.mkdir()
    main = vault / "main.md"
    outside = tmp_path / "outside.md"
    main.write_text("# Main\n\nSee [[../outside]].\n", encoding="utf-8")
    outside.write_text("# Outside\n", encoding="utf-8")

    linked = collect_linked_markdown_files([main], vault, max_depth=1)

    assert linked == [main]
    assert build_vault_index(vault)


def test_outline_note_deck_supports_explicit_proposal_mode(tmp_path: Path):
    note = tmp_path / "proposal.md"
    note.write_text(
        "\n".join(
            [
                "# Restoration Proposal",
                "",
                "## 研究假设",
                "",
                "更稳定的退化建模可以提升泛化。",
                "",
                "## 里程碑",
                "",
                "- 数据整理",
                "- 基线复现",
                "",
                "## 风险",
                "",
                "真实噪声与合成噪声分布不一致。",
            ]
        ),
        encoding="utf-8",
    )

    result = run_script(str(note), "--mode", "proposal", "--max-slides", "12", "--language", "en")

    assert "# Restoration Proposal" in result.stdout
    assert "Deck Mode: proposal" in result.stdout
    assert "Data requirements and evaluation plan" in result.stdout
    assert "Risks, mitigations, and fallback paths" in result.stdout


def test_outline_note_deck_fails_without_markdown(tmp_path: Path):
    result = subprocess.run(
        [sys.executable, "scripts/outline_note_deck.py", str(tmp_path)],
        cwd=ROOT,
        text=True,
        encoding="utf-8",
        errors="replace",
        capture_output=True,
    )

    assert result.returncode == 1
    assert "no Markdown note files found" in result.stderr


def test_build_scientific_deck_generates_nonempty_pptx(tmp_path: Path):
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "method.md").write_text(
        "# Method Note\n\n## 方法\n\nA claim.\n\n## 关键公式\n\n$$x=y$$\n\n## 实验\n\n| Metric | Value |\n| --- | --- |\n| PSNR | 30 |\n",
        encoding="utf-8",
    )
    brief = tmp_path / "brief.md"
    deck = tmp_path / "test_deck.pptx"
    run_script(str(notes), "--out", str(brief), "--title", "Method Deck", "--language", "en")

    result = run_build_script(str(brief), "--out", str(deck))

    assert f"wrote_pptx {deck}" in result.stdout
    assert deck.exists()
    assert deck.stat().st_size > 1000
    prs = Presentation(str(deck))
    assert len(prs.slides) == 15
    assert [slide_title(slide) for slide in list(prs.slides)[:4]] == [
        "Method Deck",
        "Title and research question",
        "Why this problem matters",
        "Gap in existing work",
    ]
    with ZipFile(deck) as archive:
        names = set(archive.namelist())
    assert "[Content_Types].xml" in names
    assert "ppt/presentation.xml" in names
    assert any(name.startswith("ppt/slides/slide") for name in names)


def test_build_scientific_deck_rejects_title_that_cannot_fit_fixed_box(tmp_path: Path):
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "note.md").write_text("# Note\n\nClaim.\n", encoding="utf-8")

    with pytest.raises(ValueError, match="maximum is 180"):
        build_deck(notes, tmp_path / "too-long.pptx", title="x" * 181)


def test_build_scientific_deck_respects_max_slides_from_notes_folder(tmp_path: Path):
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "paper.md").write_text(
        "# Paper Note\n\n## Problem\n\nText.\n\n## Method\n\nText.\n\n## Experiment\n\n| Metric | Value |\n| --- | --- |\n| PSNR | 30 |\n",
        encoding="utf-8",
    )
    deck = tmp_path / "limited.pptx"

    result = run_build_script(str(notes), "--out", str(deck), "--title", "Limited Deck", "--max-slides", "6", "--language", "en")

    assert "slides 6" in result.stdout
    prs = Presentation(str(deck))
    assert len(prs.slides) == 6
    assert [slide_title(slide) for slide in prs.slides] == [
        "Limited Deck",
        "Title and research question",
        "Why this problem matters",
        "Gap in existing work",
        "Limitations and open questions",
        "Appendix index",
    ]


def test_build_scientific_deck_follow_links_adds_linked_notes_to_brief(tmp_path: Path):
    vault = tmp_path / "vault"
    main = vault / "main.md"
    linked = vault / "Linked Evidence.md"
    vault.mkdir()
    main.write_text("# Main Claim\n\nSee [[Linked Evidence]].\n", encoding="utf-8")
    linked.write_text(
        "# Linked Evidence\n\n## Experiment\n\n| Metric | Value |\n| --- | --- |\n| Accuracy | 95 |\n",
        encoding="utf-8",
    )

    brief, _ = load_or_create_brief(
        main,
        title="Linked Deck",
        audience="committee review",
        max_slides=7,
        language="en",
        follow_links=True,
        vault_root=vault,
        max_depth=1,
    )

    assert "Audience: committee review" in brief
    assert "Maximum total slide count: 7 (including the title slide)" in brief
    assert "`main.md`" in brief
    assert "`Linked Evidence.md`" in brief
    assert "result/comparison table" in brief

    deck = tmp_path / "linked.pptx"
    result = run_build_script(
        str(main),
        "--out",
        str(deck),
        "--title",
        "Linked Deck",
        "--max-slides",
        "7",
        "--language",
        "en",
        "--follow-links",
        "--vault-root",
        str(vault),
        "--max-depth",
        "1",
    )

    assert "slides 7" in result.stdout
    assert len(Presentation(str(deck)).slides) == 7


def test_outline_folder_scan_skips_external_file_and_directory_symlinks(tmp_path: Path) -> None:
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "safe.md").write_text("# Safe Note\n\nIn-root evidence.\n", encoding="utf-8")
    outside_file = tmp_path / "outside.md"
    outside_file.write_text("# OUTSIDE_SECRET_FILE\n", encoding="utf-8")
    (notes / "linked.md").symlink_to(outside_file)
    outside_dir = tmp_path / "outside-dir"
    outside_dir.mkdir()
    (outside_dir / "secret.md").write_text("# OUTSIDE_SECRET_DIR\n", encoding="utf-8")
    (notes / "linked-dir").symlink_to(outside_dir, target_is_directory=True)

    result = run_script(str(notes), "--language", "en")

    assert "Safe Note" in result.stdout
    assert "OUTSIDE_SECRET_FILE" not in result.stdout
    assert "OUTSIDE_SECRET_DIR" not in result.stdout


def test_outline_explicit_markdown_file_symlink_remains_supported(tmp_path: Path) -> None:
    outside = tmp_path / "outside.md"
    outside.write_text("# Explicit Linked Note\n\nUser-selected evidence.\n", encoding="utf-8")
    link = tmp_path / "selected.md"
    link.symlink_to(outside)

    result = run_script(str(link), "--language", "en")

    assert "Explicit Linked Note" in result.stdout


def test_build_scientific_deck_folder_input_propagates_safe_scan(tmp_path: Path) -> None:
    notes = tmp_path / "notes"
    notes.mkdir()
    (notes / "safe.md").write_text("# Safe Note\n\nIn-root evidence.\n", encoding="utf-8")
    outside = tmp_path / "outside.md"
    outside.write_text("# OUTSIDE_SECRET_BUILD\n", encoding="utf-8")
    (notes / "linked.md").symlink_to(outside)

    brief, _ = load_or_create_brief(
        notes,
        title="Safe Build",
        audience="research seminar",
        max_slides=6,
        language="en",
    )

    assert "Safe Note" in brief
    assert "OUTSIDE_SECRET_BUILD" not in brief


def unsafe_output_path(tmp_path: Path, suffix: str, kind: str) -> tuple[Path, Path]:
    outside = tmp_path / "outside"
    outside.mkdir()
    if kind == "final":
        target = outside / f"sentinel{suffix}"
        target.write_bytes(b"outside sentinel\n")
        output = tmp_path / f"output{suffix}"
        output.symlink_to(target)
        return output, target
    if kind == "parent":
        parent = tmp_path / "linked-parent"
        parent.symlink_to(outside, target_is_directory=True)
        return parent / f"output{suffix}", outside / f"output{suffix}"
    ancestor = tmp_path / "linked-ancestor"
    ancestor.symlink_to(outside, target_is_directory=True)
    return ancestor / "nested" / f"output{suffix}", outside / "nested" / f"output{suffix}"


@pytest.mark.parametrize("kind", ["final", "parent", "ancestor"])
def test_outline_note_deck_rejects_output_symlink_components(tmp_path: Path, kind: str) -> None:
    note = tmp_path / "note.md"
    note.write_text("# Evidence\n\nMeasured result.\n", encoding="utf-8")
    output, outside_target = unsafe_output_path(tmp_path, ".md", kind)
    original = outside_target.read_bytes() if outside_target.exists() else None

    result = run_script_unchecked(str(note), "--out", str(output), "--language", "en")

    assert result.returncode == 1
    assert "symlink" in result.stderr.lower()
    if original is None:
        assert not outside_target.exists()
    else:
        assert outside_target.read_bytes() == original


@pytest.mark.parametrize("kind", ["final", "parent", "ancestor"])
def test_build_scientific_deck_rejects_output_symlink_components(tmp_path: Path, kind: str) -> None:
    note = tmp_path / "note.md"
    note.write_text("# Evidence\n\nMeasured result.\n", encoding="utf-8")
    output, outside_target = unsafe_output_path(tmp_path, ".pptx", kind)
    original = outside_target.read_bytes() if outside_target.exists() else None

    result = run_build_script_unchecked(
        str(note),
        "--out",
        str(output),
        "--title",
        "Safe Deck",
        "--max-slides",
        "4",
        "--language",
        "en",
    )

    assert result.returncode == 1
    assert "symlink" in result.stderr.lower()
    if original is None:
        assert not outside_target.exists()
    else:
        assert outside_target.read_bytes() == original
