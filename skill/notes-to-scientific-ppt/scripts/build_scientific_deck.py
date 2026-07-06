#!/usr/bin/env python3
"""Build an editable scientific PPTX skeleton from a deck brief or notes folder."""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path
import re
import sys
from xml.sax.saxutils import escape
from zipfile import ZIP_DEFLATED, ZipFile

try:
    from .outline_note_deck import build_brief
except ImportError:
    from outline_note_deck import build_brief


SECTION_RE = re.compile(r"^##\s+(.+?)\s*$", re.M)
SPINE_HEADINGS = {"Suggested Scientific Deck Spine", "建议科学演示主线"}
BACKLOG_HEADINGS = {"Draft Slide Backlog", "草稿幻灯片待办"}
NUMBERED_RE = re.compile(r"^\s*\d+\.\s+(.+?)\s*$")
BACKLOG_RE = re.compile(r"^-\s+\[([^\]]+)\].*?`([^`]+)`.*?Proof object:\s*(.+?)\.?\s*$")


@dataclass(frozen=True)
class SlideSpec:
    title: str
    role: str
    proof_object: str
    source: str = ""


def configure_output_encoding() -> None:
    for stream in (sys.stdout, sys.stderr):
        if hasattr(stream, "reconfigure"):
            stream.reconfigure(encoding="utf-8", errors="replace")


def is_brief(text: str) -> bool:
    return any(f"## {heading}" in text for heading in SPINE_HEADINGS | BACKLOG_HEADINGS)


def section_text(text: str, names: set[str]) -> str:
    matches = list(SECTION_RE.finditer(text))
    for index, match in enumerate(matches):
        if match.group(1).strip() not in names:
            continue
        start = match.end()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        return text[start:end]
    return ""


def parse_title(text: str, fallback: str) -> str:
    for line in text.splitlines():
        if line.startswith("# "):
            return line.removeprefix("# ").strip() or fallback
    return fallback


def infer_role(title: str) -> str:
    lower = title.lower()
    if any(token in lower for token in ("formula", "algorithm", "equation", "公式", "算法")):
        return "formula/algorithm"
    if any(token in lower for token in ("result", "experiment", "evidence", "table", "ablation", "实验", "结果", "证据", "表")):
        return "evidence/results"
    if any(token in lower for token in ("limitation", "risk", "failure", "局限", "风险", "失败")):
        return "limitations"
    if any(token in lower for token in ("appendix", "backup", "附录")):
        return "appendix"
    if any(token in lower for token in ("method", "mechanism", "pipeline", "方法", "机制")):
        return "method/mechanism"
    return "claim"


def parse_backlog(text: str) -> list[SlideSpec]:
    specs: list[SlideSpec] = []
    for line in section_text(text, BACKLOG_HEADINGS).splitlines():
        match = BACKLOG_RE.match(line.strip())
        if not match:
            continue
        role, title, proof = match.groups()
        specs.append(SlideSpec(title=title, role=role, proof_object=proof))
    return specs


def parse_spine(text: str) -> list[SlideSpec]:
    specs: list[SlideSpec] = []
    backlog = parse_backlog(text)
    backlog_by_role: dict[str, list[SlideSpec]] = {}
    for item in backlog:
        backlog_by_role.setdefault(item.role, []).append(item)

    for line in section_text(text, SPINE_HEADINGS).splitlines():
        match = NUMBERED_RE.match(line)
        if not match:
            continue
        title = match.group(1).strip()
        role = infer_role(title)
        proof = "claim with source-grounded proof object"
        for candidate_role, candidates in backlog_by_role.items():
            if role == candidate_role and candidates:
                proof = candidates[0].proof_object
                break
        specs.append(SlideSpec(title=title, role=role, proof_object=proof))
    return specs


def load_or_create_brief(input_path: Path, *, title: str | None = None, language: str = "auto", mode: str = "auto") -> tuple[str, str]:
    if input_path.is_dir():
        brief = build_brief([input_path], title, "research seminar", 18, mode, language=language)
        return brief, title or input_path.name
    if not input_path.exists():
        raise ValueError(f"input does not exist: {input_path}")
    text = input_path.read_text(encoding="utf-8", errors="replace")
    if input_path.suffix.lower() == ".md" and is_brief(text):
        return text, parse_title(text, title or input_path.stem)
    brief = build_brief([input_path], title, "research seminar", 18, mode, language=language)
    return brief, title or input_path.stem


def slide_blocks(spec: SlideSpec, slide_number: int) -> list[str]:
    if spec.role == "formula/algorithm":
        return [
            "Claim: state what this formula or algorithm proves.",
            "Formula / algorithm block: insert the key equation or pseudocode.",
            "Variable legend: define every variable and assumption.",
            f"Proof object: {spec.proof_object}.",
        ]
    if spec.role == "evidence/results":
        return [
            "Claim: state the result before showing numbers.",
            "Evidence table placeholder: task | dataset/source | metric | interpretation.",
            "What changed: explain the comparison or ablation.",
            f"Proof object: {spec.proof_object}.",
        ]
    if spec.role == "limitations":
        return [
            "Known limitation or threat to validity.",
            "Where the method or claim may fail.",
            "Missing evidence to collect before final presentation.",
            f"Proof object: {spec.proof_object}.",
        ]
    if spec.role == "appendix":
        return [
            "Appendix index.",
            "Add derivations, raw tables, extended examples, and source excerpts here.",
            "Keep dense support material out of main claim slides.",
        ]
    return [
        "Claim: rewrite this title as a falsifiable scientific claim.",
        "Evidence: cite source note, URL, figure, table, or formula.",
        f"Proof object: {spec.proof_object}.",
        "Takeaway: one sentence the audience should remember.",
    ]


def build_with_python_pptx(title: str, specs: list[SlideSpec], out: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title(slide, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
        frame = box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 38, 51)

    def add_footer(slide, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.25))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(96, 104, 116)
        p.alignment = PP_ALIGN.RIGHT

    def add_bullets(slide, bullets: list[str]) -> None:
        box = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(11.8), Inches(4.9))
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        for index, bullet in enumerate(bullets):
            p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(17)
            p.font.color.rgb = RGBColor(40, 48, 60)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(title_slide, title)
    add_bullets(title_slide, ["Editable scientific deck skeleton", "Generated from deck brief; replace placeholders with source-grounded proof objects."])
    add_footer(title_slide, "notes-to-scientific-ppt skeleton")

    for index, spec in enumerate(specs, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, spec.title)
        add_bullets(slide, slide_blocks(spec, index))
        add_footer(slide, f"Slide {index}: {spec.role}")

    prs.save(out)


def xml_text_body(lines: list[str]) -> str:
    paragraphs = []
    for line in lines:
        paragraphs.append(
            "<a:p><a:r><a:rPr lang=\"en-US\" sz=\"1800\"/>"
            f"<a:t>{escape(line)}</a:t></a:r></a:p>"
        )
    return "<p:txBody><a:bodyPr wrap=\"square\"/><a:lstStyle/>" + "".join(paragraphs) + "</p:txBody>"


def slide_xml(title: str, lines: list[str]) -> str:
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp><p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="500000" y="300000"/><a:ext cx="11500000" cy="800000"/></a:xfrm></p:spPr><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="en-US" sz="3000" b="1"/><a:t>{escape(title)}</a:t></a:r></a:p></p:txBody></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="3" name="Body"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="700000" y="1350000"/><a:ext cx="10800000" cy="4800000"/></a:xfrm></p:spPr>{xml_text_body(lines)}</p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>"""


def build_minimal_pptx(title: str, specs: list[SlideSpec], out: Path) -> None:
    slides = [(title, ["Editable scientific deck skeleton", "Generated from deck brief; replace placeholders with proof objects."])]
    slides.extend((spec.title, slide_blocks(spec, index)) for index, spec in enumerate(specs, start=1))

    overrides = [
        '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
    ]
    rels = []
    sld_ids = []
    for index in range(1, len(slides) + 1):
        overrides.append(
            f'<Override PartName="/ppt/slides/slide{index}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        )
        rels.append(
            f'<Relationship Id="rId{index}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{index}.xml"/>'
        )
        sld_ids.append(f'<p:sldId id="{255 + index}" r:id="rId{index}"/>')

    presentation_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldIdLst>{''.join(sld_ids)}</p:sldIdLst>
  <p:sldSz cx="12192000" cy="6858000" type="wide"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>"""

    content_types = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  {''.join(overrides)}
</Types>"""
    root_rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>"""
    presentation_rels = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">{''.join(rels)}</Relationships>"""

    with ZipFile(out, "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", root_rels)
        archive.writestr("ppt/presentation.xml", presentation_xml)
        archive.writestr("ppt/_rels/presentation.xml.rels", presentation_rels)
        for index, (slide_title, lines) in enumerate(slides, start=1):
            archive.writestr(f"ppt/slides/slide{index}.xml", slide_xml(slide_title, lines))


def build_deck(input_path: Path, out: Path, *, title: str | None = None, language: str = "auto", mode: str = "auto") -> int:
    brief, fallback_title = load_or_create_brief(input_path, title=title, language=language, mode=mode)
    deck_title = parse_title(brief, title or fallback_title)
    specs = parse_spine(brief)
    if not specs:
        specs = parse_backlog(brief)
    if not specs:
        specs = [SlideSpec("Core claim from notes", "claim", "claim-and-evidence text slide")]
    if not any(spec.role == "limitations" for spec in specs):
        specs.append(SlideSpec("Limitations and open questions", "limitations", "limitation/failure-case slide"))
    if not any(spec.role == "appendix" for spec in specs):
        specs.append(SlideSpec("Appendix index", "appendix", "appendix navigation"))

    out.parent.mkdir(parents=True, exist_ok=True)
    try:
        build_with_python_pptx(deck_title, specs, out)
    except ModuleNotFoundError:
        build_minimal_pptx(deck_title, specs, out)
    return len(specs) + 1


def main() -> int:
    configure_output_encoding()
    parser = argparse.ArgumentParser(description="Build an editable scientific PPTX skeleton from a deck brief or notes folder.")
    parser.add_argument("input", type=Path, help="Deck brief Markdown, note file, or notes folder")
    parser.add_argument("--out", required=True, type=Path, help="Output .pptx path")
    parser.add_argument("--title", help="Deck title override when building directly from notes")
    parser.add_argument("--language", choices=["zh", "en", "auto"], default="auto")
    parser.add_argument("--mode", choices=["auto", "paper-reading", "proposal", "progress-report", "teaching", "defense"], default="auto")
    args = parser.parse_args()

    if args.out.suffix.lower() != ".pptx":
        parser.error("--out must end with .pptx")
    try:
        slide_count = build_deck(args.input, args.out, title=args.title, language=args.language, mode=args.mode)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    print(f"wrote_pptx {args.out}")
    print(f"slides {slide_count}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
