#!/usr/bin/env python3
"""Reopen a generated PPTX, verify slide count, and optionally render it."""

from __future__ import annotations

import argparse
import math
import os
from pathlib import Path
import re
import shutil
import stat
import subprocess
import sys
import tempfile
from zipfile import ZipFile


SLIDE_RE = re.compile(r"^ppt/slides/slide([0-9]+)\.xml$")
SLIDE_PNG_RE = re.compile(r"^slide-([0-9]+)\.png$", re.I)


class RenderPublishRecoveryError(RuntimeError):
    """Raised when a render publish cannot restore the previous output set."""


def package_slide_count(path: Path) -> int:
    with ZipFile(path) as archive:
        if archive.testzip() is not None:
            raise ValueError(f"PPTX contains a corrupt ZIP member: {archive.testzip()}")
        names = set(archive.namelist())
        required = {"[Content_Types].xml", "ppt/presentation.xml"}
        missing = sorted(required - names)
        if missing:
            raise ValueError(f"PPTX is missing required parts: {', '.join(missing)}")
        slide_numbers = sorted(int(match.group(1)) for name in names if (match := SLIDE_RE.match(name)))
    if slide_numbers != list(range(1, len(slide_numbers) + 1)):
        raise ValueError(f"PPTX slide parts are not contiguous: {slide_numbers}")
    return len(slide_numbers)


def reopen_presentation(path: Path):
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise RuntimeError("python-pptx is required for the reopen gate") from exc
    return Presentation(str(path))


def slide_count(presentation) -> int:
    return len(presentation.slides)


def first_slide_text(presentation) -> str:
    if not presentation.slides:
        return ""
    return "\n".join(
        shape.text.strip()
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def relative_beneath(root: Path, path: Path) -> Path:
    root = root.resolve()
    candidate = Path(os.path.abspath(path.expanduser()))
    try:
        return candidate.relative_to(root)
    except ValueError as exc:
        raise ValueError(f"render output must remain beside the PPTX: {path}") from exc


def ensure_safe_directory(root: Path, path: Path, *, create: bool) -> Path:
    root = root.resolve()
    relative = relative_beneath(root, path)
    current = root
    for component in relative.parts:
        current = current / component
        try:
            mode = current.lstat().st_mode
        except FileNotFoundError:
            if not create:
                break
            current.mkdir()
            mode = current.lstat().st_mode
        if stat.S_ISLNK(mode):
            raise ValueError(f"render directory contains symlink component: {current}")
        if not stat.S_ISDIR(mode):
            raise ValueError(f"render directory component is not a directory: {current}")
    return root / relative


def ensure_safe_generated_file(root: Path, path: Path) -> Path:
    root = root.resolve()
    relative = relative_beneath(root, path)
    candidate = root / relative
    ensure_safe_directory(root, candidate.parent, create=False)
    try:
        mode = candidate.lstat().st_mode
    except FileNotFoundError:
        return candidate
    if stat.S_ISLNK(mode):
        raise ValueError(f"generated render path is a symlink: {candidate}")
    if not stat.S_ISREG(mode):
        raise ValueError(f"generated render path is not a regular file: {candidate}")
    return candidate


def directory_identity(path: Path) -> tuple[int, int]:
    mode = path.lstat()
    if stat.S_ISLNK(mode.st_mode) or not stat.S_ISDIR(mode.st_mode):
        raise ValueError(f"render directory is not a regular directory: {path}")
    return mode.st_dev, mode.st_ino


def revalidate_render_directory(
    root: Path,
    output_dir: Path,
    expected_identity: tuple[int, int],
) -> Path:
    safe_output = ensure_safe_directory(root, output_dir, create=False)
    if directory_identity(safe_output) != expected_identity:
        raise ValueError(f"render directory identity changed during rendering: {safe_output}")
    return safe_output


def indexed_slide_targets(root: Path, output_dir: Path) -> dict[int, Path]:
    indexed: dict[int, Path] = {}
    for path in output_dir.iterdir():
        match = SLIDE_PNG_RE.fullmatch(path.name)
        if match is None:
            continue
        safe_path = ensure_safe_generated_file(root, path)
        page = int(match.group(1))
        if page in indexed:
            raise ValueError(
                f"generated slide previews contain duplicate page {page}: "
                f"{indexed[page].name}, {path.name}"
            )
        indexed[page] = safe_path
    return indexed


def validate_rendered_slide_set(
    root: Path,
    output_dir: Path,
    pages: int,
) -> list[Path]:
    indexed = indexed_slide_targets(root, output_dir)
    expected = set(range(1, pages + 1))
    if set(indexed) != expected:
        raise RuntimeError(
            "rendered PNG page set does not match PDF pages: "
            f"actual={sorted(indexed)} expected={sorted(expected)}"
        )
    return [indexed[page] for page in range(1, pages + 1)]


def open_directory_fd(path: Path) -> int | None:
    """Hold a directory identity during publish on platforms with dir-fd APIs."""

    if os.name == "nt":
        return None
    flags = os.O_RDONLY
    for name in ("O_DIRECTORY", "O_CLOEXEC", "O_NOFOLLOW"):
        flags |= getattr(os, name, 0)
    return os.open(path, flags)


def entry_identity(directory_fd: int | None, path: Path) -> tuple[int, int] | None:
    try:
        if directory_fd is None:
            status = path.lstat()
        else:
            status = os.stat(path.name, dir_fd=directory_fd, follow_symlinks=False)
    except FileNotFoundError:
        return None
    return status.st_dev, status.st_ino


def replace_entry(
    source: Path,
    target: Path,
    source_fd: int | None,
    target_fd: int | None,
) -> None:
    if source_fd is None or target_fd is None:
        os.replace(source, target)
        return
    os.replace(
        source.name,
        target.name,
        src_dir_fd=source_fd,
        dst_dir_fd=target_fd,
    )


def publish_render_outputs(
    render_root: Path,
    output_dir: Path,
    render_identity: tuple[int, int],
    staged_pdf: Path,
    staged_slides: list[Path],
    backup_dir: Path,
) -> tuple[Path, list[Path]]:
    """Replace the owned render set, rolling back in-process interruptions.

    This transaction does not promise atomic recovery from process termination
    or machine failure. An incomplete in-process rollback keeps its private
    recovery directory and reports that path instead of deleting evidence.
    """

    output_dir = revalidate_render_directory(render_root, output_dir, render_identity)
    final_pdf = ensure_safe_generated_file(
        render_root,
        output_dir / staged_pdf.name,
    )
    existing_slides = indexed_slide_targets(render_root, output_dir)
    existing = [path for path in (final_pdf, *existing_slides.values()) if path.exists()]
    staging_dir = staged_pdf.parent
    expected_directory_identities = (
        render_identity,
        directory_identity(staging_dir),
        directory_identity(backup_dir),
    )
    output_fd: int | None = None
    staging_fd: int | None = None
    backup_fd: int | None = None
    try:
        output_fd = open_directory_fd(output_dir)
        staging_fd = open_directory_fd(staging_dir)
        backup_fd = open_directory_fd(backup_dir)
        for descriptor, directory, expected_identity in zip(
            (output_fd, staging_fd, backup_fd),
            (output_dir, staging_dir, backup_dir),
            expected_directory_identities,
        ):
            if descriptor is None:
                continue
            opened = os.fstat(descriptor)
            if (opened.st_dev, opened.st_ino) != expected_identity:
                raise ValueError(
                    f"render transaction directory identity changed: {directory}"
                )
    except BaseException:
        for descriptor in (backup_fd, staging_fd, output_fd):
            if descriptor is not None:
                os.close(descriptor)
        raise

    backups: list[tuple[Path, Path, tuple[int, int]]] = []
    published: list[tuple[Path, Path, tuple[int, int]]] = []
    try:
        for path in existing:
            output_dir = revalidate_render_directory(
                render_root,
                output_dir,
                render_identity,
            )
            backup = backup_dir / path.name
            identity = entry_identity(output_fd, path)
            if identity is None:
                raise ValueError(f"render output changed before backup: {path}")
            backups.append((backup, path, identity))
            replace_entry(path, backup, output_fd, backup_fd)

        for staged in (staged_pdf, *staged_slides):
            output_dir = revalidate_render_directory(
                render_root,
                output_dir,
                render_identity,
            )
            target = ensure_safe_generated_file(
                render_root,
                output_dir / staged.name,
            )
            identity = entry_identity(staging_fd, staged)
            if identity is None:
                raise ValueError(f"staged render output changed before publish: {staged}")
            published.append((target, staged, identity))
            replace_entry(staged, target, staging_fd, output_fd)

        final_pdf = ensure_safe_generated_file(
            render_root,
            output_dir / staged_pdf.name,
        )
        final_slides = validate_rendered_slide_set(
            render_root,
            output_dir,
            len(staged_slides),
        )
        return final_pdf, final_slides
    except BaseException as original:
        rollback_errors: list[str] = []
        for target, staged, identity in reversed(published):
            try:
                staged_identity = entry_identity(staging_fd, staged)
                target_identity = entry_identity(output_fd, target)
                if target_identity == identity and staged_identity is None:
                    replace_entry(target, staged, output_fd, staging_fd)
                elif staged_identity != identity:
                    rollback_errors.append(
                        f"published {target.name} has conflicting recovery state"
                    )
            except BaseException as exc:
                rollback_errors.append(f"recover published {target.name}: {exc}")
        for backup, target, identity in reversed(backups):
            try:
                backup_identity = entry_identity(backup_fd, backup)
                target_identity = entry_identity(output_fd, target)
                if backup_identity == identity and target_identity is None:
                    replace_entry(backup, target, backup_fd, output_fd)
                elif target_identity != identity:
                    rollback_errors.append(
                        f"backup {target.name} has conflicting recovery state"
                    )
            except BaseException as exc:
                rollback_errors.append(f"restore {target.name}: {exc}")
        if rollback_errors:
            raise RenderPublishRecoveryError(
                "render publish failed; previous outputs need recovery from "
                f"{backup_dir.parent}: "
                + "; ".join(rollback_errors)
            ) from original
        raise
    finally:
        for descriptor in (backup_fd, staging_fd, output_fd):
            if descriptor is not None:
                os.close(descriptor)


def render_pptx(path: Path, output_dir: Path, expected_slides: int | None, require_render: bool) -> None:
    converter = shutil.which("soffice") or shutil.which("libreoffice")
    renderer = shutil.which("pdftoppm")
    pdfinfo = shutil.which("pdfinfo")
    if converter is None or renderer is None or pdfinfo is None:
        message = "LibreOffice/Poppler render tools are unavailable"
        if require_render:
            raise RuntimeError(f"MANUAL_REVIEW_REQUIRED: {message}")
        print(f"MANUAL_REVIEW_REQUIRED: {message}")
        return

    selected_root = Path(os.path.abspath(path.parent.expanduser()))
    selected_output = Path(os.path.abspath(output_dir.expanduser()))
    try:
        output_relative = selected_output.relative_to(selected_root)
    except ValueError as exc:
        raise ValueError(f"render output must remain beside the PPTX: {output_dir}") from exc
    render_root = path.parent.expanduser().resolve()
    output_dir = ensure_safe_directory(render_root, render_root / output_relative, create=True)
    render_identity = directory_identity(output_dir)
    ensure_safe_generated_file(render_root, output_dir / f"{path.stem}.pdf")
    indexed_slide_targets(render_root, output_dir)
    transaction_dir = Path(
        tempfile.mkdtemp(
            prefix=f".{output_dir.name}-stage-",
            dir=render_root,
        )
    )
    try:
        staging_dir = transaction_dir / "rendered"
        backup_dir = transaction_dir / "backup"
        staging_dir.mkdir()
        backup_dir.mkdir()
        staged_pdf = staging_dir / f"{path.stem}.pdf"
        with tempfile.TemporaryDirectory(prefix="scientific-deck-lo-") as profile_name:
            profile_dir = Path(profile_name)
            subprocess.run(
                [
                    converter,
                    "--headless",
                    f"-env:UserInstallation={profile_dir.as_uri()}",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(staging_dir),
                    str(path),
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
                encoding="utf-8",
                errors="replace",
            )
        output_dir = revalidate_render_directory(
            render_root,
            output_dir,
            render_identity,
        )
        staged_pdf = ensure_safe_generated_file(render_root, staged_pdf)
        if not staged_pdf.is_file():
            raise RuntimeError(f"LibreOffice did not create {staged_pdf}")
        info = subprocess.run(
            [pdfinfo, str(staged_pdf)],
            check=True,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
        )
        output_dir = revalidate_render_directory(
            render_root,
            output_dir,
            render_identity,
        )
        staged_pdf = ensure_safe_generated_file(render_root, staged_pdf)
        pages_match = re.search(r"^Pages:\s*(\d+)\s*$", info.stdout, re.M)
        if pages_match is None:
            raise RuntimeError("pdfinfo did not report a page count")
        pages = int(pages_match.group(1))
        if expected_slides is not None and pages != expected_slides:
            raise RuntimeError(
                f"rendered PDF has {pages} pages; expected {expected_slides}"
            )
        subprocess.run(
            [renderer, "-png", str(staged_pdf), str(staging_dir / "slide")],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
        )
        output_dir = revalidate_render_directory(
            render_root,
            output_dir,
            render_identity,
        )
        staged_pdf = ensure_safe_generated_file(render_root, staged_pdf)
        staged_slides = validate_rendered_slide_set(
            render_root,
            staging_dir,
            pages,
        )
        pdf, _slides = publish_render_outputs(
            render_root,
            output_dir,
            render_identity,
            staged_pdf,
            staged_slides,
            backup_dir,
        )
    except RenderPublishRecoveryError:
        raise
    except BaseException:
        shutil.rmtree(transaction_dir, ignore_errors=True)
        raise
    else:
        try:
            shutil.rmtree(transaction_dir)
        except OSError as exc:
            raise RuntimeError(
                "render outputs were committed, but transaction cleanup failed at "
                f"{transaction_dir}: {exc}"
            ) from exc
    print(f"render ok pdf={pdf} pages={pages}")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--expected-slides", type=int)
    parser.add_argument("--expected-title")
    parser.add_argument("--expected-width-inches", type=float)
    parser.add_argument("--expected-height-inches", type=float)
    parser.add_argument("--render", action="store_true", help="render to PDF and PNG previews when tools are available")
    parser.add_argument(
        "--require-render",
        action="store_true",
        help="require rendering (implies --render) and fail when render tools are unavailable or rendering fails",
    )
    args = parser.parse_args(argv)
    try:
        package_slides = package_slide_count(args.pptx)
        presentation = reopen_presentation(args.pptx)
        reopened_slides = slide_count(presentation)
        if package_slides != reopened_slides:
            raise RuntimeError(f"package has {package_slides} slides but python-pptx reopened {reopened_slides}")
        if args.expected_slides is not None and reopened_slides != args.expected_slides:
            raise RuntimeError(f"reopened PPTX has {reopened_slides} slides; expected {args.expected_slides}")
        if args.expected_title and not first_slide_text(presentation).startswith(args.expected_title):
            raise RuntimeError(f"first-slide title does not start with expected title {args.expected_title!r}")
        for label, expected, actual in (
            ("width", args.expected_width_inches, presentation.slide_width / 914400),
            ("height", args.expected_height_inches, presentation.slide_height / 914400),
        ):
            if expected is not None and not math.isclose(actual, expected, rel_tol=0, abs_tol=0.01):
                raise RuntimeError(f"slide {label} is {actual:.3f} inches; expected {expected:.3f}")
        print(f"pptx_reopen ok slides={reopened_slides}")
        if args.render or args.require_render:
            render_pptx(args.pptx, args.pptx.parent / f"{args.pptx.stem}-render", args.expected_slides or reopened_slides, args.require_render)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
