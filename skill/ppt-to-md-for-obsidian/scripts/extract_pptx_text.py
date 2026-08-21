#!/usr/bin/env python3
"""Extract ordered text from a PPTX file into Markdown.

This script is intentionally conservative: it extracts visible text, table
cells, and speaker notes when python-pptx exposes them. The output is raw
material for rewriting, not final Obsidian notes.
"""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
import re
import sys
from zipfile import BadZipFile, ZipFile
import xml.etree.ElementTree as ET

try:
    from .safe_io import InputTooLargeError, ensure_safe_input_file, read_bytes_no_follow, safe_write_text
except ImportError:
    from safe_io import InputTooLargeError, ensure_safe_input_file, read_bytes_no_follow, safe_write_text


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
}
SLIDE_XML_RE = re.compile(r"ppt/slides/slide(\d+)\.xml$")
REQUIRED_PPTX_PARTS = {"[Content_Types].xml", "ppt/presentation.xml"}
MAX_PPTX_INPUT_BYTES = 256 * 1024 * 1024
MAX_PPTX_ZIP_MEMBERS = 20_000
MAX_PPTX_MEMBER_BYTES = 128 * 1024 * 1024
MAX_PPTX_TOTAL_UNCOMPRESSED_BYTES = 1024 * 1024 * 1024
MAX_PPTX_COMPRESSION_RATIO = 1000


class PptxExtractionError(ValueError):
    """Stable user-facing failure for an unreadable PPTX input."""

    def __init__(self, path: Path, reason: str) -> None:
        self.path = path
        self.reason = reason
        super().__init__(f"{path}: {reason}")


@dataclass
class ShapeRecord:
    top: int
    left: int
    text: str
    kind: str = "text"


@dataclass(frozen=True)
class PptxExtractionResult:
    markdown: str
    backend: str
    partial: bool
    slide_count: int
    blank_slides: int
    media_objects: int


def _stable_pptx_bytes(path: Path) -> tuple[Path, bytes]:
    if path.suffix.casefold() != ".pptx":
        raise PptxExtractionError(path, "input must be a .pptx file")
    try:
        path = ensure_safe_input_file(path)
    except (OSError, ValueError) as exc:
        reason = str(exc)
        if "does not exist" in reason:
            reason = "file does not exist"
        elif "not a regular file" in reason:
            reason = "input is not a regular file"
        raise PptxExtractionError(path, reason) from exc

    try:
        return path, read_bytes_no_follow(path, max_bytes=MAX_PPTX_INPUT_BYTES)
    except InputTooLargeError as exc:
        raise PptxExtractionError(
            path, f"PPTX input exceeds {MAX_PPTX_INPUT_BYTES} byte safety limit"
        ) from exc


def _validate_zip_budget(path: Path, archive: ZipFile) -> None:
    members = archive.infolist()
    if len(members) > MAX_PPTX_ZIP_MEMBERS:
        raise PptxExtractionError(path, "PPTX ZIP package has too many members")
    names: set[str] = set()
    total = 0
    for member in members:
        if member.filename in names:
            raise PptxExtractionError(path, "PPTX ZIP package contains duplicate members")
        names.add(member.filename)
        if member.file_size > MAX_PPTX_MEMBER_BYTES:
            raise PptxExtractionError(path, "PPTX ZIP member exceeds the uncompressed size limit")
        total += member.file_size
        if total > MAX_PPTX_TOTAL_UNCOMPRESSED_BYTES:
            raise PptxExtractionError(path, "PPTX ZIP package exceeds the total uncompressed size limit")
        if member.file_size and (
            member.compress_size == 0
            or member.file_size > member.compress_size * MAX_PPTX_COMPRESSION_RATIO
        ):
            raise PptxExtractionError(path, "PPTX ZIP member exceeds the compression-ratio limit")
    if any(member.flag_bits & 0x1 for member in members):
        raise PptxExtractionError(path, "encrypted PPTX packages are not supported")


def _validated_pptx_payload(path: Path) -> tuple[Path, bytes]:
    """Read one stable PPTX snapshot and validate its OOXML boundary."""

    path, payload = _stable_pptx_bytes(path)
    try:
        with ZipFile(BytesIO(payload)) as archive:
            _validate_zip_budget(path, archive)
            members = archive.infolist()
            if archive.testzip() is not None:
                raise PptxExtractionError(path, "PPTX ZIP package has a corrupt member")
            names = {member.filename for member in members}
            if not REQUIRED_PPTX_PARTS <= names:
                raise PptxExtractionError(
                    path,
                    "PPTX package is missing required OOXML parts",
                )
            try:
                for name in REQUIRED_PPTX_PARTS:
                    ET.fromstring(archive.read(name))
            except (ET.ParseError, UnicodeError):
                raise PptxExtractionError(
                    path,
                    "PPTX package contains malformed required XML",
                ) from None
    except PptxExtractionError:
        raise
    except (BadZipFile, EOFError):
        raise PptxExtractionError(path, "invalid PPTX ZIP package") from None
    except OSError as exc:
        raise PptxExtractionError(path, "PPTX package cannot be read") from exc
    return path, payload


def validate_pptx_input(path: Path) -> Path:
    """Validate one stable ZIP/OOXML snapshot before selecting a backend."""

    validated, _payload = _validated_pptx_payload(path)
    return validated


def extraction_header(
    path: Path,
    *,
    backend: str,
    partial: bool,
    slide_count: int,
    blank_slides: int,
    media_objects: int,
) -> list[str]:
    lines = [
        f"# Extracted PPTX Text: {path.name}",
        "",
        f"- Backend: `{backend}`",
        f"- Partial fallback: {str(partial).lower()}",
        f"- Slides: {slide_count}",
        f"- Slides without visible text: {blank_slides}",
        f"- Media objects: {media_objects}",
        "",
    ]
    if partial:
        lines.extend(
            [
                "Warning: ZIP/XML fallback is partial; speaker notes, media meaning, and some relationships may be missing.",
                "Use OCR or manual slide inspection before claiming complete source coverage.",
                "",
            ]
        )
    if blank_slides or media_objects:
        lines.extend(
            [
                "Coverage warning: slides without visible text or with media require visual inspection; text extraction alone is not complete visual/OCR coverage.",
                "",
            ]
        )
    return lines


def position(shape) -> tuple[int, int]:
    top = int(getattr(shape, "top", 0) or 0)
    left = int(getattr(shape, "left", 0) or 0)
    return top, left


def iter_shape_text(shape):
    if getattr(shape, "has_text_frame", False) and shape.text:
        yield shape.text

    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                yield " | ".join(cells)

    if getattr(shape, "shapes", None):
        for subshape in shape.shapes:
            yield from iter_shape_text(subshape)


def iter_shape_records(shape, include_media_placeholders: bool = True):
    top, left = position(shape)
    emitted_text = False
    for text in iter_shape_text(shape):
        text = text.strip()
        if text:
            emitted_text = True
            yield ShapeRecord(top=top, left=left, text=text, kind="text")

    if emitted_text or not include_media_placeholders:
        return

    shape_type = str(getattr(shape, "shape_type", "")).lower()
    name = getattr(shape, "name", "")
    if "picture" in shape_type:
        yield ShapeRecord(top=top, left=left, text=f"[Image placeholder: {name or 'picture'}]", kind="image")
    elif getattr(shape, "has_chart", False):
        yield ShapeRecord(top=top, left=left, text=f"[Chart placeholder: {name or 'chart'}]", kind="chart")


def shape_media_count(shape) -> int:
    shape_type = str(getattr(shape, "shape_type", "")).lower()
    count = int("picture" in shape_type or bool(getattr(shape, "has_chart", False)))
    for subshape in getattr(shape, "shapes", ()) or ():
        count += shape_media_count(subshape)
    return count


def extract_notes(slide):
    try:
        notes_slide = getattr(slide, "notes_slide", None)
    except Exception:
        notes_slide = None
    if notes_slide is None:
        return []
    chunks = []
    for shape in notes_slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text:
            text = shape.text.strip()
            if text and text.lower() != "click to add notes":
                chunks.append(text)
    return chunks


def xml_shape_position(shape: ET.Element) -> tuple[int, int]:
    off = shape.find(".//a:xfrm/a:off", NS)
    if off is None:
        return 0, 0
    return int(off.attrib.get("y", "0") or 0), int(off.attrib.get("x", "0") or 0)


def xml_shape_text(shape: ET.Element) -> str:
    paragraphs = []
    for paragraph in shape.findall(".//a:p", NS):
        text = "".join(node.text or "" for node in paragraph.findall(".//a:t", NS)).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def sorted_slide_xml_names(archive: ZipFile) -> list[str]:
    matches = []
    for name in archive.namelist():
        match = SLIDE_XML_RE.fullmatch(name)
        if match:
            matches.append((int(match.group(1)), name))
    return [name for _, name in sorted(matches)]


def _extract_pptx_with_zip_payload(
    path: Path, payload: bytes, include_slide_title: bool = True
) -> PptxExtractionResult:
    with ZipFile(BytesIO(payload)) as archive:
        _validate_zip_budget(path, archive)
        slide_names = sorted_slide_xml_names(archive)
        rendered_slides: list[str] = []
        blank_slides = 0
        media_objects = 0
        for idx, slide_name in enumerate(slide_names, start=1):
            root = ET.fromstring(archive.read(slide_name))
            records = []
            seen = set()
            for shape in root.findall(".//p:sp", NS) + root.findall(".//p:graphicFrame", NS):
                text = xml_shape_text(shape).strip()
                if not text:
                    continue
                top, left = xml_shape_position(shape)
                key = (text, top, left)
                if key in seen:
                    continue
                records.append(ShapeRecord(top=top, left=left, text=text, kind="text"))
                seen.add(key)
            records.sort(key=lambda record: (record.top, record.left))
            if not records:
                blank_slides += 1
            media_objects += len(root.findall(".//p:pic", NS))
            media_objects += len(root.findall(".//c:chart", NS))

            title = slide_title(records) if include_slide_title else None
            if title:
                rendered_slides.append(f"## Slide {idx}: {title}")
            else:
                rendered_slides.append(f"## Slide {idx}")
            rendered_slides.append("")

            if records:
                for record in records:
                    for line in record.text.splitlines():
                        line = line.strip()
                        if line:
                            rendered_slides.append(f"- {line}")
            else:
                rendered_slides.append("- [No visible text extracted]")
            rendered_slides.append("")

    out = extraction_header(
        path,
        backend="zip-xml-fallback",
        partial=True,
        slide_count=len(slide_names),
        blank_slides=blank_slides,
        media_objects=media_objects,
    )
    out.extend(rendered_slides)
    return PptxExtractionResult(
        markdown="\n".join(out).rstrip() + "\n",
        backend="zip-xml-fallback",
        partial=True,
        slide_count=len(slide_names),
        blank_slides=blank_slides,
        media_objects=media_objects,
    )


def extract_pptx_with_zip_result(path: Path, include_slide_title: bool = True) -> PptxExtractionResult:
    path, payload = _validated_pptx_payload(path)
    try:
        return _extract_pptx_with_zip_payload(
            path, payload, include_slide_title=include_slide_title
        )
    except PptxExtractionError:
        raise
    except (BadZipFile, EOFError, OSError, ET.ParseError) as exc:
        raise PptxExtractionError(path, "PPTX package could not be parsed") from exc


def extract_pptx_with_zip(path: Path, include_slide_title: bool = True) -> str:
    return extract_pptx_with_zip_result(path, include_slide_title=include_slide_title).markdown


def slide_title(records: list[ShapeRecord]) -> str | None:
    for record in records:
        if record.kind != "text":
            continue
        for line in record.text.splitlines():
            line = line.strip()
            if line:
                return line
    return None


def extract_pptx_result(
    path: Path,
    include_media_placeholders: bool = True,
    include_slide_title: bool = True,
) -> PptxExtractionResult:
    original_path = path
    path, payload = _validated_pptx_payload(path)
    try:
        from pptx import Presentation
    except ImportError:
        try:
            return _extract_pptx_with_zip_payload(
                path, payload, include_slide_title=include_slide_title
            )
        except Exception as exc:
            raise PptxExtractionError(original_path, "PPTX package could not be parsed") from exc

    try:
        prs = Presentation(BytesIO(payload))
        rendered_slides: list[str] = []
        blank_slides = 0
        media_objects = 0

        for idx, slide in enumerate(prs.slides, start=1):
            records = []
            seen = set()
            for shape in slide.shapes:
                for record in iter_shape_records(shape, include_media_placeholders=include_media_placeholders):
                    key = (record.text, record.top, record.left)
                    if key not in seen:
                        records.append(record)
                        seen.add(key)
            records.sort(key=lambda record: (record.top, record.left))
            if not any(record.kind == "text" for record in records):
                blank_slides += 1
            media_objects += sum(shape_media_count(shape) for shape in slide.shapes)

            title = slide_title(records) if include_slide_title else None
            if title:
                rendered_slides.append(f"## Slide {idx}: {title}")
            else:
                rendered_slides.append(f"## Slide {idx}")
            rendered_slides.append("")

            if records:
                for record in records:
                    chunk = record.text
                    for line in chunk.splitlines():
                        line = line.strip()
                        if line:
                            rendered_slides.append(f"- {line}")
            else:
                rendered_slides.append("- [No visible text extracted]")

            notes = extract_notes(slide)
            if notes:
                rendered_slides.append("")
                rendered_slides.append("### Speaker Notes")
                for note in notes:
                    for line in note.splitlines():
                        line = line.strip()
                        if line:
                            rendered_slides.append(f"- {line}")

            rendered_slides.append("")

        out = extraction_header(
            path,
            backend="python-pptx",
            partial=False,
            slide_count=len(prs.slides),
            blank_slides=blank_slides,
            media_objects=media_objects,
        )
        out.extend(rendered_slides)
        return PptxExtractionResult(
            markdown="\n".join(out).rstrip() + "\n",
            backend="python-pptx",
            partial=False,
            slide_count=len(prs.slides),
            blank_slides=blank_slides,
            media_objects=media_objects,
        )
    except Exception as exc:
        raise PptxExtractionError(original_path, "PPTX package could not be parsed") from exc


def extract_pptx(path: Path, include_media_placeholders: bool = True, include_slide_title: bool = True) -> str:
    return extract_pptx_result(
        path,
        include_media_placeholders=include_media_placeholders,
        include_slide_title=include_slide_title,
    ).markdown


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract PPTX text into Markdown.")
    parser.add_argument("pptx", type=Path, help="Path to a .pptx file")
    parser.add_argument("--out", type=Path, help="Output Markdown path")
    parser.add_argument("--no-media-placeholders", action="store_true", help="Do not emit image/chart placeholders.")
    parser.add_argument("--no-slide-title", action="store_true", help="Do not add detected slide titles to headings.")
    args = parser.parse_args()

    try:
        md = extract_pptx(
            args.pptx,
            include_media_placeholders=not args.no_media_placeholders,
            include_slide_title=not args.no_slide_title,
        )
    except PptxExtractionError as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    if args.out:
        try:
            safe_write_text(args.out, md)
        except (OSError, ValueError) as exc:
            print(f"ERROR: {exc}", file=sys.stderr)
            return 1
    else:
        print(md, end="")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
