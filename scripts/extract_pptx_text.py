#!/usr/bin/env python3
"""Extract ordered text from a PPTX file into Markdown.

This script is intentionally conservative: it extracts visible text, table
cells, and speaker notes when python-pptx exposes them. The output is raw
material for rewriting, not final Obsidian notes.
"""

from __future__ import annotations

import argparse
from pathlib import Path
import sys


def iter_shape_text(shape):
    if getattr(shape, "has_text_frame", False) and shape.text:
        yield shape.text

    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                yield " | ".join(cells)

    if getattr(shape, "shapes", None):
        for subshape in shape.shapes:
            yield from iter_shape_text(subshape)


def extract_notes(slide):
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return []
    chunks = []
    for shape in notes_slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text:
            text = shape.text.strip()
            if text and text.lower() != "click to add notes":
                chunks.append(text)
    return chunks


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit(
            "Missing dependency: python-pptx. Install it with "
            "`python3 -m pip install python-pptx` in the active environment."
        )

    prs = Presentation(str(path))
    out = [f"# Extracted PPTX Text: {path.name}", ""]

    for idx, slide in enumerate(prs.slides, start=1):
        out.append(f"## Slide {idx}")
        out.append("")

        chunks = []
        for shape in slide.shapes:
            for text in iter_shape_text(shape):
                text = text.strip()
                if text and text not in chunks:
                    chunks.append(text)

        if chunks:
            for chunk in chunks:
                for line in chunk.splitlines():
                    line = line.strip()
                    if line:
                        out.append(f"- {line}")
        else:
            out.append("- [No visible text extracted]")

        notes = extract_notes(slide)
        if notes:
            out.append("")
            out.append("### Speaker Notes")
            for note in notes:
                for line in note.splitlines():
                    line = line.strip()
                    if line:
                        out.append(f"- {line}")

        out.append("")

    return "\n".join(out).rstrip() + "\n"


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract PPTX text into Markdown.")
    parser.add_argument("pptx", type=Path, help="Path to a .pptx file")
    parser.add_argument("--out", type=Path, help="Output Markdown path")
    args = parser.parse_args()

    if not args.pptx.exists():
        parser.error(f"file does not exist: {args.pptx}")
    if args.pptx.suffix.lower() != ".pptx":
        parser.error("input must be a .pptx file")

    md = extract_pptx(args.pptx)
    if args.out:
        args.out.write_text(md, encoding="utf-8")
    else:
        print(md, end="")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
